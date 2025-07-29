"""Office format processors (DOCX, XLSX, PPTX)."""

import logging
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Union, Tuple

from doc2mark.core.base import (
    BaseProcessor,
    DocumentFormat,
    DocumentMetadata,
    ProcessedDocument,
    ProcessingError
)
from doc2mark.ocr.base import BaseOCR

logger = logging.getLogger(__name__)


class OfficeProcessor(BaseProcessor):
    """Processor for modern Office formats (DOCX, XLSX, PPTX)."""

    def __init__(self, ocr: Optional[BaseOCR] = None):
        """Initialize Office processor.
        
        Args:
            ocr: OCR provider for image extraction
        """
        self.ocr = ocr
        self._docx = None
        self._openpyxl = None
        self._pptx = None

    @property
    def python_docx(self):
        """Lazy load python-docx."""
        if self._docx is None:
            try:
                import docx
                self._docx = docx
            except ImportError:
                raise ImportError(
                    "python-docx is not installed. "
                    "Install it with: pip install python-docx"
                )
        return self._docx

    @property
    def openpyxl(self):
        """Lazy load openpyxl."""
        if self._openpyxl is None:
            try:
                import openpyxl
                self._openpyxl = openpyxl
            except ImportError:
                raise ImportError(
                    "openpyxl is not installed. "
                    "Install it with: pip install openpyxl"
                )
        return self._openpyxl

    @property
    def python_pptx(self):
        """Lazy load python-pptx."""
        if self._pptx is None:
            try:
                import pptx
                self._pptx = pptx
            except ImportError:
                raise ImportError(
                    "python-pptx is not installed. "
                    "Install it with: pip install python-pptx"
                )
        return self._pptx

    def can_process(self, file_path: Union[str, Path]) -> bool:
        """Check if this processor can handle the file."""
        file_path = Path(file_path)
        extension = file_path.suffix.lower().lstrip('.')
        return extension in ['docx', 'xlsx', 'pptx']

    def process(
            self,
            file_path: Union[str, Path],
            **kwargs
    ) -> ProcessedDocument:
        """Process Office document."""
        file_path = Path(file_path)
        extension = file_path.suffix.lower().lstrip('.')

        # Get file size
        file_size = file_path.stat().st_size

        # Process based on format
        if extension == 'docx':
            content, metadata = self._process_docx(file_path, **kwargs)
            doc_format = DocumentFormat.DOCX
        elif extension == 'xlsx':
            content, metadata = self._process_xlsx(file_path, **kwargs)
            doc_format = DocumentFormat.XLSX
        elif extension == 'pptx':
            content, metadata = self._process_pptx(file_path, **kwargs)
            doc_format = DocumentFormat.PPTX
        else:
            raise ProcessingError(f"Unsupported Office format: {extension}")

        # Build metadata
        doc_metadata = DocumentMetadata(
            filename=file_path.name,
            format=doc_format,
            size_bytes=file_size,
            **metadata
        )

        return ProcessedDocument(
            content=content,
            metadata=doc_metadata
        )

    def _process_docx(self, file_path: Path, **kwargs) -> Tuple[str, dict]:
        """Process DOCX document."""
        try:
            doc = self.python_docx.Document(str(file_path))

            # Extract text content
            markdown_parts = []

            # Process paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    # Check style for headings
                    if para.style.name.startswith('Heading'):
                        level = int(para.style.name[-1]) if para.style.name[-1].isdigit() else 1
                        markdown_parts.append(f"{'#' * level} {para.text}")
                    else:
                        # Handle text formatting
                        text = self._format_docx_paragraph(para)
                        markdown_parts.append(text)
                    markdown_parts.append("")

            # Process tables
            for table in doc.tables:
                table_md = self._convert_docx_table_to_markdown(table)
                markdown_parts.append(table_md)
                markdown_parts.append("")

            # Extract images if requested
            images = []
            if kwargs.get('extract_images', False) and self.ocr:
                images = self._extract_docx_images(file_path)

            # Metadata
            metadata = {
                'page_count': len(doc.element.xpath('//w:sectPr')),
                'word_count': sum(len(para.text.split()) for para in doc.paragraphs),
                'author': doc.core_properties.author,
                'title': doc.core_properties.title,
                'creation_date': str(doc.core_properties.created) if doc.core_properties.created else None,
                'modification_date': str(doc.core_properties.modified) if doc.core_properties.modified else None,
            }

            return '\n'.join(markdown_parts), metadata

        except Exception as e:
            logger.error(f"Failed to process DOCX: {e}")
            raise ProcessingError(f"DOCX processing failed: {str(e)}")

    def _process_xlsx(self, file_path: Path, **kwargs) -> Tuple[str, dict]:
        """Process XLSX document."""
        try:
            wb = self.openpyxl.load_workbook(str(file_path), data_only=True)

            markdown_parts = []
            total_cells = 0

            # Process each sheet
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]

                if sheet.max_row > 0 and sheet.max_column > 0:
                    markdown_parts.append(f"## {sheet_name}")
                    markdown_parts.append("")

                    # Convert sheet to markdown table
                    table_data = []
                    for row in sheet.iter_rows(values_only=True):
                        # Filter out completely empty rows
                        if any(cell is not None for cell in row):
                            table_data.append([str(cell) if cell is not None else "" for cell in row])
                            total_cells += len(row)

                    if table_data:
                        table_md = self._convert_list_to_markdown_table(table_data)
                        markdown_parts.append(table_md)
                        markdown_parts.append("")

            # Metadata
            metadata = {
                'page_count': len(wb.sheetnames),
                'sheet_names': wb.sheetnames,
                'total_cells': total_cells
            }

            return '\n'.join(markdown_parts), metadata

        except Exception as e:
            logger.error(f"Failed to process XLSX: {e}")
            raise ProcessingError(f"XLSX processing failed: {str(e)}")

    def _process_pptx(self, file_path: Path, **kwargs) -> Tuple[str, dict]:
        """Process PPTX document."""
        try:
            prs = self.python_pptx.Presentation(str(file_path))

            markdown_parts = []
            image_count = 0

            # Process each slide
            for i, slide in enumerate(prs.slides, 1):
                markdown_parts.append(f"## Slide {i}")
                markdown_parts.append("")

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Handle title shapes
                        if shape == slide.shapes.title:
                            markdown_parts.append(f"### {shape.text}")
                        else:
                            markdown_parts.append(shape.text)
                        markdown_parts.append("")

                    # Count images
                    if shape.shape_type == 13:  # Picture
                        image_count += 1
                        if kwargs.get('extract_images', False) and self.ocr:
                            # TODO: Extract and process image
                            pass

                # Process tables
                for shape in slide.shapes:
                    if shape.has_table:
                        table_md = self._convert_pptx_table_to_markdown(shape.table)
                        markdown_parts.append(table_md)
                        markdown_parts.append("")

            # Metadata
            metadata = {
                'page_count': len(prs.slides),
                'slide_count': len(prs.slides),
                'image_count': image_count,
                'title': prs.core_properties.title,
                'author': prs.core_properties.author,
            }

            return '\n'.join(markdown_parts), metadata

        except Exception as e:
            logger.error(f"Failed to process PPTX: {e}")
            raise ProcessingError(f"PPTX processing failed: {str(e)}")

    def _format_docx_paragraph(self, paragraph) -> str:
        """Format DOCX paragraph with inline formatting."""
        if not paragraph.runs:
            return paragraph.text

        formatted_text = []
        for run in paragraph.runs:
            text = run.text
            if run.bold:
                text = f"**{text}**"
            if run.italic:
                text = f"*{text}*"
            formatted_text.append(text)

        return ''.join(formatted_text)

    def _convert_docx_table_to_markdown(self, table) -> str:
        """Convert DOCX table to markdown."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        return self._convert_list_to_markdown_table(rows)

    def _convert_pptx_table_to_markdown(self, table) -> str:
        """Convert PPTX table to markdown."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        return self._convert_list_to_markdown_table(rows)

    def _convert_list_to_markdown_table(self, data: List[List[str]]) -> str:
        """Convert list of lists to markdown table."""
        if not data:
            return ""

        # Determine column widths
        col_widths = [0] * len(data[0])
        for row in data:
            for i, cell in enumerate(row):
                col_widths[i] = max(col_widths[i], len(cell))

        # Build table
        lines = []

        # Header
        header = "| " + " | ".join(cell.ljust(col_widths[i]) for i, cell in enumerate(data[0])) + " |"
        lines.append(header)

        # Separator
        separator = "|" + "|".join("-" * (w + 2) for w in col_widths) + "|"
        lines.append(separator)

        # Data rows
        for row in data[1:]:
            row_str = "| " + " | ".join(cell.ljust(col_widths[i]) for i, cell in enumerate(row)) + " |"
            lines.append(row_str)

        return '\n'.join(lines)

    def _extract_docx_images(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract images from DOCX file using batch OCR."""
        images = []

        try:
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # Collect all images for batch processing
                image_batch = []
                image_info = []

                # Find image files in the media folder
                for file_info in zip_file.filelist:
                    if file_info.filename.startswith('word/media/'):
                        image_data = zip_file.read(file_info.filename)
                        image_batch.append(image_data)
                        image_info.append({
                            'filename': file_info.filename,
                            'size': len(image_data)
                        })

                # Batch process OCR if we have images and OCR is available
                if image_batch and self.ocr:
                    try:
                        # Prepare OCR kwargs with language configuration if available
                        ocr_kwargs = {}
                        if hasattr(self.ocr, 'config') and self.ocr.config and self.ocr.config.language:
                            ocr_kwargs['language'] = self.ocr.config.language
                            logger.info(
                                f"🌍 Passing language configuration to Office batch OCR: {self.ocr.config.language}")

                        logger.info(f"🚀 Processing {len(image_batch)} DOCX images with batch OCR")
                        ocr_results = self.ocr.batch_process_images(image_batch, **ocr_kwargs)

                        # Combine results
                        for info, ocr_result in zip(image_info, ocr_results):
                            text = ocr_result.text if hasattr(ocr_result, 'text') else str(ocr_result)
                            info['text'] = text
                            images.append(info)

                        logger.info(f"✅ DOCX batch OCR completed successfully")

                    except Exception as e:
                        logger.warning(f"Failed to batch OCR DOCX images: {e}")
                        # Add images without OCR text
                        for info in image_info:
                            info['text'] = ''
                            images.append(info)
                else:
                    # No OCR available, add images without text
                    for info in image_info:
                        info['text'] = ''
                        images.append(info)

        except Exception as e:
            logger.warning(f"Failed to extract images from DOCX: {e}")

        return images
