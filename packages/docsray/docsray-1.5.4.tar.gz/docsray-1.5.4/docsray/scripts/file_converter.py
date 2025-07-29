#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
docsray/scripts/file_converter.py
Convert various file formats to PDF for processing
"""

import os
import sys
from pathlib import Path
from typing import Optional, Tuple

import pypandoc

from PIL import Image

import pdfkit

import markdown

from llama_index.readers.file import HWPReader

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from PIL import Image
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch    

from bs4 import BeautifulSoup


from typing import List, Tuple
from pathlib import Path
from docsray.config import CACHE_DIR
import pandas as pd

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.utils import ImageReader
import os

def _save_text_images_to_pdf_korean(text: str, image_paths: List[str], output_file: Path) -> bool:
    """Save text and images to PDF with Korean font support"""
    try:          
        # Register Korean fonts
        # Try to find Korean fonts on the system
        korean_fonts = [
            # Windows
            "C:/Windows/Fonts/malgun.ttf",  # Malgun Gothic
            "C:/Windows/Fonts/gulim.ttc",   # Gulim
            "C:/Windows/Fonts/batang.ttc",  # Batang
            # macOS
            "/System/Library/Fonts/AppleSDGothicNeo.ttc",
            "/Library/Fonts/NanumGothic.ttf",
            # Linux
            "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
            "/usr/share/fonts/truetype/fonts-nanum/NanumGothic.ttf",
        ]
        
        font_registered = False
        font_name = "NanumGothic"
        
        for font_path in korean_fonts:
            if os.path.exists(font_path):
                try:
                    pdfmetrics.registerFont(TTFont(font_name, font_path))
                    font_registered = True
                    break
                except:
                    continue
        
        if not font_registered:
            # Fallback: download and use a free Korean font
            import urllib.request
            font_url = "https://github.com/google/fonts/raw/main/ofl/nanumgothic/NanumGothic-Regular.ttf"
            font_path = output_file.parent / "NanumGothic.ttf"
            urllib.request.urlretrieve(font_url, font_path)
            pdfmetrics.registerFont(TTFont(font_name, str(font_path)))
        
        # Create PDF
        c = canvas.Canvas(str(output_file), pagesize=A4)
        width, height = A4
        
        # Set font with Korean support
        c.setFont(font_name, 12)
        
        # Add text
        y = height - 50
        for line in text.split('\n'):
            if y < 50:  # New page if needed
                c.showPage()
                c.setFont(font_name, 12)
                y = height - 50
            
            # Handle long lines
            if len(line) > 80:
                words = line.split()
                current_line = ""
                for word in words:
                    if len(current_line + word) < 80:
                        current_line += word + " "
                    else:
                        c.drawString(50, y, current_line.strip())
                        y -= 20
                        current_line = word + " "
                if current_line:
                    c.drawString(50, y, current_line.strip())
                    y -= 20
            else:
                c.drawString(50, y, line)
                y -= 20
        
        # Add images
        for img_path in image_paths:
            if os.path.exists(img_path):
                try:
                    c.showPage()
                    img = ImageReader(img_path)
                    c.drawImage(img, 50, 50, width=width-100, height=height-100, preserveAspectRatio=True)
                except:
                    pass
        
        c.save()
        return True
        
    except Exception as e:
        print(f"PDF creation error: {e}", file=sys.stderr)
        return False
            
class FileConverter:
    """Convert various file formats to PDF"""
    
    SUPPORTED_FORMATS = {
        # Office documents
        '.docx': 'Microsoft Word',
        '.doc': 'Microsoft Word (Legacy)',
        '.xlsx': 'Microsoft Excel',
        '.xls': 'Microsoft Excel (Legacy)',
        '.pdf': 'PDF Document',
        '.pptx': 'Microsoft PowerPoint',
        '.ppt': 'Microsoft PowerPoint (Legacy)',
        '.odt': 'OpenDocument Text',
        '.ods': 'OpenDocument Spreadsheet',
        '.odp': 'OpenDocument Presentation',
        '.hwp': 'Hancom Word Processor',
        '.hwpx': 'Hancom Word Processor (OOXML)',

        # Text formats
        '.txt': 'Plain Text',
        
        # Image formats
        '.jpg': 'JPEG Image',
        '.jpeg': 'JPEG Image',
        '.png': 'PNG Image',
        '.gif': 'GIF Image',
        '.bmp': 'Bitmap Image',
        '.tiff': 'TIFF Image',
        '.tif': 'TIFF Image',
        '.webp': 'WebP Image',
    }
    
    def __init__(self, output_dir: Optional[Path] = None):
        """
        Initialize converter with optional output directory
        
        Args:
            output_dir: Directory to save converted PDFs (default: temp directory)
        """
        self.output_dir = output_dir or CACHE_DIR
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
    @classmethod
    def get_supported_formats(cls) -> dict:
        """Get dictionary of supported formats"""
        return cls.SUPPORTED_FORMATS.copy()
    
    @classmethod
    def is_supported(cls, file_path: str) -> bool:
        """Check if file format is supported"""
        ext = Path(file_path).suffix.lower()
        return ext in cls.SUPPORTED_FORMATS or ext == '.pdf'

    def convert_to_pdf(self, input_path: str, output_path: Optional[str] = None) -> Tuple[bool, str]:
        """
        Convert file to PDF
        
        Args:
            input_path: Path to input file
            output_path: Optional output path (default: auto-generated)
            
        Returns:
            Tuple of (success: bool, output_path_or_error: str)
        """
        input_file = Path(input_path)
        
        # Check if file exists
        if not input_file.exists():
            return False, f"File not found: {input_path}"
        
        # Check if already PDF
        if input_file.suffix.lower() == '.pdf':
            return True, str(input_file)
        
        # Check if format is supported
        file_ext = input_file.suffix.lower()
        if file_ext not in self.SUPPORTED_FORMATS:
            return False, f"Unsupported format: {file_ext}"
        
        # Generate output path if not provided
        if output_path is None:
            output_path = self.output_dir / f"{input_file.stem}_converted.pdf"
        else:
            output_path = Path(output_path)
        
        # Select conversion method based on file type
        print(f"Converting {self.SUPPORTED_FORMATS[file_ext]} file to PDF...", file=sys.stderr)
        
        # Office documents
        if file_ext == '.docx':
            return self._convert_docx_to_pdf(input_file, output_path)
        elif file_ext == '.doc':
            return self._convert_doc_to_pdf(input_file, output_path)
        elif file_ext in ['.xlsx', '.xls']:
            return self._convert_excel_to_pdf(input_file, output_path)
        elif file_ext in ['.pptx', '.ppt']:
            return self._convert_ppt_to_pdf(input_file, output_path)
        elif file_ext == '.hwp':
            return self._convert_hwp_to_pdf(input_file, output_path)
        elif file_ext == '.hwpx':
            return self._convert_hwpx_to_pdf(input_file, output_path)
        
        elif file_ext in ['.pdf']:
            # If already PDF, just return the path
            return True, str(input_file)

        # Text formats
        elif file_ext == '.txt':
            return self._convert_text_to_pdf(input_file, output_path)
        
        # Image formats
        elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp']:
            return self._convert_image_to_pdf(input_file, output_path)
        
        # Fallback: try pandoc for anything else
        else:
            return self._convert_with_pandoc(input_file, output_path)
        
    # ------------------------------------------------------------------
    # HWPX → PDF  (zip + XML)  — text + images via BeautifulSoup + ReportLab
    # ------------------------------------------------------------------
    def _convert_hwpx_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert HWPX (Hangul OOXML) to PDF."""
        from zipfile import ZipFile
        texts, img_paths = [], []
        tmp_dir = output_file.with_suffix("")  # temp folder alongside output
        tmp_dir.mkdir(exist_ok=True)

        try:
            with ZipFile(input_file) as zf:
                # 1) Collect text + image ids
                for name in zf.namelist():
                    if name.startswith("Contents/") and name.endswith(".xml"):
                        # Read with proper encoding
                        raw_content = zf.read(name)
                        
                        # Try different encodings for Korean text
                        for encoding in ['utf-8', 'utf-8-sig', 'cp949', 'euc-kr']:
                            try:
                                content = raw_content.decode(encoding)
                                break
                            except UnicodeDecodeError:
                                continue
                        else:
                            # If all encodings fail, use utf-8 with error handling
                            content = raw_content.decode('utf-8', errors='replace')
                        
                        soup = BeautifulSoup(content, "xml")
                        
                        # HWPX uses different tag structure
                        # Look for text in various possible tags
                        text_tags = ['t', 'text', 'hp:t', 'hp:run', 'w:t']
                        for tag_name in text_tags:
                            for text_elem in soup.find_all(tag_name):
                                text = text_elem.get_text(strip=True)
                                if text:
                                    texts.append(text)
                        
                        # Also try to find paragraphs with different naming conventions
                        para_tags = ['p', 'para', 'hp:p', 'hp:para', 'w:p']
                        for tag_name in para_tags:
                            for para in soup.find_all(tag_name):
                                para_text = para.get_text(" ", strip=True)
                                if para_text and para_text not in texts:
                                    texts.append(para_text)
                        
                        # Find images with various possible tag/attribute names
                        img_tags = ['pic', 'hp:pic', 'image', 'hp:image']
                        img_attrs = ['binFile', 'hp:binFile', 'r:embed', 'href']
                        
                        for img_tag in img_tags:
                            for pic in soup.find_all(img_tag):
                                for attr in img_attrs:
                                    img_id = pic.get(attr)
                                    if img_id:
                                        img_paths.append(img_id)

                # Debug: Print what we found
                print(f"Found {len(texts)} text elements", file=sys.stderr)
                print(f"Found {len(img_paths)} image references", file=sys.stderr)
                
                # 2) Extract BinData images
                extracted_imgs = []
                for img_id in img_paths:
                    for possible_path in [f"BinData/{img_id}", f"Contents/BinData/{img_id}", img_id]:
                        try:
                            data = zf.read(possible_path)
                            # Detect image type
                            if data[:8] == b"\x89PNG\r\n\x1a\n":
                                ext = ".png"
                            elif data[:2] == b"\xff\xd8":
                                ext = ".jpg"
                            else:
                                ext = ".bin"  # unknown format
                            
                            out_path = tmp_dir / f"{input_file.stem}_{img_id}{ext}"
                            out_path.write_bytes(data)
                            extracted_imgs.append(str(out_path))
                            break
                        except KeyError:
                            continue

                # 3) Synthesize PDF with proper encoding
                combined_text = "\n".join(texts) if texts else "No text content found"
                
                # Ensure _save_text_images_to_pdf handles Korean text properly
                success = _save_text_images_to_pdf_korean(combined_text, extracted_imgs, output_file)
                if success:
                    return True, str(output_file)
                return False, "PDF synthesis failed"

        except Exception as e:
            return False, f"HWPX conversion error: {e}"
        finally:
            # Cleanup temp directory
            try:
                import shutil
                shutil.rmtree(tmp_dir)
            except:
                pass

    def _convert_hwp_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert HWP to PDF"""
        try:
            reader = HWPReader()
            # llama_index HWPReader returns a list of Document objects
            documents = reader.load_data(file=str(input_file))
            
            if documents and len(documents) > 0:
                # Combine text from all documents
                full_text = ""
                all_images = []
                
                for doc in documents:
                    # llama_index Document objects have a 'text' attribute
                    full_text += doc.text + "\n\n"
                    
                    # Extract images from metadata
                    # Metadata structure may vary depending on HWPReader implementation
                    if hasattr(doc, 'metadata') and doc.metadata:
                        images = doc.metadata.get("images", [])
                        all_images.extend(images)
                
                # Save text and images to PDF with Korean font support
                if _save_text_images_to_pdf_korean(full_text.strip(), all_images, output_file):
                    return True, str(output_file)
                else:
                    return False, "Failed to save PDF"
            else:
                return False, "Unable to read HWP document"
                
        except ImportError as e:
            return False, f"{e} not installed"
        except Exception as e:
            error_msg = f"Error during HWP conversion: {str(e)}"
            print(f"[Error] {error_msg}", file=sys.stderr)
            return False, error_msg

    def _convert_doc_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert DOC (legacy format) to PDF"""
        # First try: Use python-docx2txt (no external dependencies)
        try:
            import docx2txt
            
            # Extract text
            text = docx2txt.process(str(input_file))
            
            # Extract images
            temp_dir = output_file.parent / f"temp_{input_file.stem}"
            temp_dir.mkdir(exist_ok=True)
            
            # docx2txt extracts images to a directory
            images_dir = temp_dir / "images"
            images_dir.mkdir(exist_ok=True)
            
            # Extract with images
            text_with_images = docx2txt.process(str(input_file), str(images_dir))
            
            # Find extracted images
            image_files = []
            if images_dir.exists():
                for img_file in images_dir.iterdir():
                    if img_file.suffix.lower() in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
                        image_files.append(str(img_file))
            
            # Convert to PDF
            if _save_text_images_to_pdf_korean(text, image_files, output_file):
                # Cleanup
                import shutil
                if temp_dir.exists():
                    shutil.rmtree(temp_dir)
                return True, str(output_file)
            else:
                return False, "Failed to save PDF"
                
        except ImportError:
            # docx2txt not installed, try alternative methods
            pass
        except Exception as e:
            print(f"docx2txt failed: {e}", file=sys.stderr)
        
        # Second try: Use python-docx (may work for some .doc files)
        try:
            from docx import Document
            
            doc = Document(str(input_file))
            full_text = ""
            
            for paragraph in doc.paragraphs:
                full_text += paragraph.text + "\n"
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = "\t".join([cell.text for cell in row.cells])
                    full_text += row_text + "\n"
                full_text += "\n"
            
            if full_text.strip():
                if _save_text_images_to_pdf_korean(full_text, [], output_file):
                    return True, str(output_file)
        except Exception as e:
            print(f"python-docx failed: {e}", file=sys.stderr)
        
        # Third try: Use pandoc if available
        try:
            # Check if pandoc is available
            pypandoc.get_pandoc_version()
            
            # Try to convert
            text = pypandoc.convert_file(str(input_file), 'plain')
            
            if text and _save_text_images_to_pdf_korean(text, [], output_file):
                return True, str(output_file)
        except Exception as e:
            print(f"Pandoc failed: {e}", file=sys.stderr)
        
        # All methods failed
        return False, (
            "Failed to convert .doc file. This is a legacy format that requires special handling.\n\n"
            "Please try one of the following:\n"
            "1. Save the file as .docx format in Microsoft Word and upload again\n"
            "2. Install docx2txt: pip install docx2txt\n"
            "3. Install python-docx: pip install python-docx\n"
            "4. Convert to PDF manually and upload the PDF\n\n"
            "Note: .docx files are recommended over .doc for better compatibility."
        )
    
    def _convert_docx_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert DOCX to PDF - try multiple methods"""
        # First try: python-docx (most reliable for .docx)
        try:
            from docx import Document
            
            doc = Document(str(input_file))
            full_text = ""
            temp_dir = output_file.parent / f"temp_{input_file.stem}"
            temp_dir.mkdir(exist_ok=True)
            image_files = []
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                full_text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = "\t".join([cell.text for cell in row.cells])
                    full_text += row_text + "\n"
                full_text += "\n"
            
            # Extract images
            try:
                img_counter = 0
                for rel in doc.part.rels.values():
                    if "image" in rel.target_ref:
                        img_counter += 1
                        img_data = rel.target_part.blob
                        img_path = temp_dir / f"image_{img_counter}.png"
                        with open(img_path, 'wb') as f:
                            f.write(img_data)
                        image_files.append(str(img_path))
            except Exception as e:
                print(f"Image extraction warning: {e}", file=sys.stderr)
            
            # Convert to PDF
            if _save_text_images_to_pdf_korean(full_text, image_files, output_file):
                # Cleanup
                import shutil
                if temp_dir.exists():
                    shutil.rmtree(temp_dir)
                return True, str(output_file)
                
        except ImportError:
            print("python-docx not installed, trying pandoc...", file=sys.stderr)
        except Exception as e:
            print(f"python-docx failed: {e}, trying pandoc...", file=sys.stderr)
        
        # Second try: pypandoc (if python-docx fails)
        try:
            import subprocess
            
            # Ensure pandoc is downloaded
            pypandoc.ensure_pandoc_installed()
            
            # Create temp directory
            temp_dir = output_file.parent / f"temp_{input_file.stem}"
            temp_dir.mkdir(exist_ok=True)
            
            # Extract text using pypandoc
            try:
                full_text = pypandoc.convert_file(str(input_file), 'plain')
            except Exception as e:
                return False, f"Pandoc text extraction failed: {str(e)}"
            
            # Extract images - still need subprocess for this
            # Get pandoc path from pypandoc
            pandoc_path = pypandoc.get_pandoc_path()
            
            subprocess.run([
                pandoc_path, str(input_file),
                '-t', 'markdown',
                '-o', str(temp_dir / 'dummy.md'),
                '--extract-media', str(temp_dir)
            ], capture_output=True)
            
            # Find extracted images
            image_files = []
            media_dir = temp_dir / "media"
            if media_dir.exists():
                for img_file in media_dir.rglob('*'):
                    if img_file.suffix.lower() in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
                        image_files.append(str(img_file))
            
            # Use Korean PDF function
            if _save_text_images_to_pdf_korean(full_text, image_files, output_file):
                # Cleanup
                import shutil
                if temp_dir.exists():
                    shutil.rmtree(temp_dir)
                return True, str(output_file)
            else:
                return False, "Failed to save PDF"
                
        except Exception as e:
            return False, f"DOCX conversion error: {str(e)}.\nPlease install python-docx: pip install python-docx"

    def _convert_excel_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert Excel to PDF using openpyxl"""
        try:
            import openpyxl
            
            # Load Excel file
            wb = openpyxl.load_workbook(input_file, data_only=True)
            
            # Extract all text content
            full_text = ""
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                full_text += f"=== Sheet: {sheet_name} ===\n\n"
                
                # Extract data from sheet
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else '' for cell in row])
                    if row_text.strip():
                        full_text += row_text + "\n"
                
                full_text += "\n\n"
            
            images = []

            if _save_text_images_to_pdf_korean(full_text.strip(), images, output_file):
                return True, str(output_file)
            else:
                return False, "Failed to save PDF"
                
        except ImportError:
            return False, "Please install openpyxl: pip install openpyxl"
        except Exception as e:
            return False, f"Excel conversion error: {str(e)}"

    def _convert_ppt_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert PowerPoint to PDF using python-pptx"""
        try:
            from pptx import Presentation
            from PIL import Image
            import io
            
            # Load PowerPoint file
            prs = Presentation(str(input_file))
            
            # Extract text and images
            full_text = ""
            image_paths = []
            temp_dir = output_file.parent / f"temp_{input_file.stem}"
            temp_dir.mkdir(exist_ok=True)
            
            for slide_num, slide in enumerate(prs.slides):
                full_text += f"=== Slide {slide_num + 1} ===\n\n"
                
                # Extract text from slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        full_text += shape.text.strip() + "\n\n"
                    
                    # Extract images
                    if shape.shape_type == 13:  # Picture
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            img = Image.open(io.BytesIO(image_bytes))
                            
                            # Save image
                            img_path = temp_dir / f"slide_{slide_num}_img_{shape.shape_id}.png"
                            img.save(img_path)
                            image_paths.append(str(img_path))
                        except:
                            pass
                
                # Extract tables
                for shape in slide.shapes:
                    if shape.has_table:
                        full_text += "[Table]\n"
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join([cell.text for cell in row.cells])
                            full_text += row_text + "\n"
                        full_text += "\n"
                
                full_text += "\n"
            
            if _save_text_images_to_pdf_korean(full_text.strip(), image_paths, output_file):
                # Cleanup temp directory
                import shutil
                shutil.rmtree(temp_dir)
                return True, str(output_file)
            else:
                return False, "Failed to save PDF"
                
        except ImportError:
            return False, "Please install python-pptx and pillow: pip install python-pptx pillow"
        except Exception as e:
            return False, f"PowerPoint conversion error: {str(e)}"
            
    def _convert_text_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert plain text to PDF"""
        try:
            # Read text file
            with open(input_file, 'r', encoding='utf-8') as f:
                text = f.read()

            try:
                doc = SimpleDocTemplate(str(output_file), pagesize=A4)
                styles = getSampleStyleSheet()
                story = []
                
                # Title
                title_style = ParagraphStyle(
                    'CustomTitle',
                    parent=styles['Heading1'],
                    fontSize=24,
                    spaceAfter=30
                )
                story.append(Paragraph(input_file.name, title_style))
                story.append(Spacer(1, 0.2*inch))
                
                # Content
                text_style = ParagraphStyle(
                    'CustomText',
                    parent=styles['Normal'],
                    fontSize=11,
                    leading=14
                )
                
                # Split text into paragraphs
                for paragraph in text.split('\n\n'):
                    if paragraph.strip():
                        # Escape special characters for reportlab
                        safe_text = paragraph.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        story.append(Paragraph(safe_text, text_style))
                        story.append(Spacer(1, 0.1*inch))
                
                doc.build(story)
                return True, str(output_file)
                
            except Exception as e:
                print(f"Reportlab conversion failed: {e}", file=sys.stderr)
            
            # Second try: Direct PDF creation using pure Python
            try:
                # Simple PDF creation without external dependencies
                pdf_content = self._create_simple_pdf(text, input_file.name)
                with open(output_file, 'wb') as f:
                    f.write(pdf_content)
                return True, str(output_file)
            except Exception as e:
                print(f"Simple PDF creation failed: {e}", file=sys.stderr)
            
            # Third try: Create HTML and convert
            html_content = f"""
            <html>
            <head>
                <meta charset="utf-8">
                <style>
                    body {{
                        font-family: Arial, sans-serif;
                        margin: 40px;
                        line-height: 1.6;
                        white-space: pre-wrap;
                    }}
                    h1 {{
                        color: #333;
                        border-bottom: 2px solid #333;
                        padding-bottom: 10px;
                    }}
                    pre {{
                        background: #f4f4f4;
                        padding: 15px;
                        border-radius: 5px;
                        overflow-x: auto;
                    }}
                </style>
            </head>
            <body>
                <h1>{input_file.name}</h1>
                <pre>{text}</pre>
            </body>
            </html>
            """
            
            # Try various HTML to PDF converters
            converters_tried = []
            # Try pdfkit
            try:
                pdfkit.from_string(html_content, str(output_file))
                return True, str(output_file)
            except Exception as e:
                converters_tried.append("pdfkit")
            
            try:
                temp_html = output_file.with_suffix('.html')
                with open(temp_html, 'w', encoding='utf-8') as f:
                    f.write(html_content)
                
                pypandoc.convert_file(str(temp_html), 'pdf', outputfile=str(output_file))
                temp_html.unlink()
                return True, str(output_file)
            except Exception as e:
                converters_tried.append("pandoc")
                    
        except Exception as e:
            print(f"Text to PDF conversion error: {e}", file=sys.stderr)
        
        # Provide helpful error message
        error_msg = "Text conversion failed. "
        if converters_tried:
            error_msg += f"Tried: {', '.join(converters_tried)}. "
        else:
            error_msg += "No PDF converters available. Install: pip install reportlab"
        
        return False, error_msg

    def _create_simple_pdf(self, text: str, title: str) -> bytes:
        """Create a very simple PDF without external dependencies"""
        # This is a minimal PDF creator - for production use reportlab is recommended
        lines = text.split('\n')
        
        # Basic PDF structure
        pdf = b"%PDF-1.4\n"
        
        # Catalog and Pages
        pdf += b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
        pdf += b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
        
        # Page
        pdf += b"3 0 obj\n<< /Type /Page /Parent 2 0 R /Resources 4 0 R /MediaBox [0 0 612 792] /Contents 5 0 R >>\nendobj\n"
        
        # Resources
        pdf += b"4 0 obj\n<< /Font << /F1 << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> >> >>\nendobj\n"
        
        # Content stream
        content = f"BT /F1 12 Tf 50 750 Td ({title}) Tj ET\n"
        y_pos = 720
        for line in lines[:50]:  # Limit to first 50 lines for simplicity
            if line.strip():
                # Escape special characters
                safe_line = line.replace('\\', '\\\\').replace('(', '\\(').replace(')', '\\)')
                content += f"BT /F1 10 Tf 50 {y_pos} Td ({safe_line[:80]}) Tj ET\n"
                y_pos -= 15
                if y_pos < 50:
                    break
        
        content_bytes = content.encode('latin-1', errors='replace')
        pdf += f"5 0 obj\n<< /Length {len(content_bytes)} >>\nstream\n".encode()
        pdf += content_bytes
        pdf += b"\nendstream\nendobj\n"
        
        # xref table
        xref_pos = len(pdf)
        pdf += b"xref\n0 6\n"
        pdf += b"0000000000 65535 f \n"
        pdf += b"0000000009 00000 n \n"
        pdf += b"0000000058 00000 n \n"
        pdf += b"0000000115 00000 n \n"
        pdf += b"0000000229 00000 n \n"
        pdf += b"0000000328 00000 n \n"
        
        # Trailer
        pdf += b"trailer\n<< /Size 6 /Root 1 0 R >>\n"
        pdf += f"startxref\n{xref_pos}\n".encode()
        pdf += b"%%EOF"
        
        return pdf
    
    def _convert_image_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert image to PDF"""       
        try:
            # Open image
            img = Image.open(input_file)
            
            # Convert to RGB if necessary
            if img.mode in ('RGBA', 'LA'):
                # Create white background
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'RGBA':
                    background.paste(img, mask=img.split()[-1])
                else:
                    background.paste(img)
                img = background
            elif img.mode not in ('RGB', 'L'):
                img = img.convert('RGB')
            
            # Save as PDF
            img.save(str(output_file), 'PDF', resolution=100.0)
            return True, str(output_file)
            
        except Exception as e:
            return False, f"Image conversion failed: {e}"
    
    def _convert_with_pandoc_simple(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert documents using Pandoc without LaTeX"""
        try:
            import subprocess
            
            # Create temp directory for extracted images
            temp_dir = output_file.parent / f"temp_{input_file.stem}"
            temp_dir.mkdir(exist_ok=True)
            
            # Extract text
            result = subprocess.run([
                'pandoc', str(input_file), '-t', 'plain'
            ], capture_output=True, text=True)
            
            if result.returncode != 0:
                return False, "Pandoc text extraction failed"
            
            full_text = result.stdout
            
            # Extract images
            subprocess.run([
                'pandoc', str(input_file),
                '-t', 'markdown',
                '-o', str(temp_dir / 'dummy.md'),
                '--extract-media', str(temp_dir)
            ], capture_output=True)
            
            # Find extracted images
            image_files = []
            media_dir = temp_dir / "media"
            if media_dir.exists():
                for img_file in media_dir.rglob('*'):
                    if img_file.suffix.lower() in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
                        image_files.append(str(img_file))
            
            # 한글 지원 PDF 함수 사용
            if _save_text_images_to_pdf_korean(full_text, image_files, output_file):
                # Cleanup
                import shutil
                shutil.rmtree(temp_dir)
                return True, str(output_file)
            else:
                return False, "Failed to save PDF"
                
        except Exception as e:
            return False, f"Conversion error: {str(e)}"

def convert_file_to_pdf(input_path: str, output_dir: Optional[str] = None) -> Tuple[bool, str]:
    """
    Convenience function to convert a file to PDF
    
    Args:
        input_path: Path to input file
        output_dir: Optional output directory
        
    Returns:
        Tuple of (success: bool, output_path_or_error: str)
    """
    converter = FileConverter(Path(output_dir) if output_dir else None)
    return converter.convert_to_pdf(input_path)


if __name__ == "__main__":
    # Test conversion
    import argparse
    
    parser = argparse.ArgumentParser(description="Convert files to PDF")
    parser.add_argument("input_file", help="Input file path")
    parser.add_argument("-o", "--output", help="Output PDF path")
    parser.add_argument("-d", "--output-dir", help="Output directory")
    
    args = parser.parse_args()
    
    converter = FileConverter(Path(args.output_dir) if args.output_dir else None)
    success, result = converter.convert_to_pdf(args.input_file, args.output)
    
    if success:
        print(f"✅ Converted successfully: {result}, file=sys.stderr")
    else:
        print(f"❌ Conversion failed: {result}", file=sys.stderr)
        sys.exit(1)
