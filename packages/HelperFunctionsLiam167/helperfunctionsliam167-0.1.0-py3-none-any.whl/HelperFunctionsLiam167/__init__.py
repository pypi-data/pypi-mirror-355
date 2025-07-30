# @title 🎨 Plotly Config
import json
import math
import os
import pickle
from datetime import timedelta
import pandas as pd
import numpy as np
import duckdb
import datetime
from google.colab import data_table
from scipy import stats
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from matplotlib.ticker import FuncFormatter, FixedLocator
import seaborn as sns
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import plotly.io as pio
from sklearn.preprocessing import MinMaxScaler
from sklearn.model_selection import train_test_split
from sklearn.ensemble import RandomForestRegressor
from sklearn.metrics import mean_squared_error, r2_score
import ipywidgets as widgets
from IPython.display import display, HTML, clear_output
import builtins
from sklearn.ensemble import RandomForestClassifier
from xgboost import XGBClassifier
import xlsxwriter
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from google.colab import files, output
from pptx.dml.color import RGBColor
import io
from google.colab import drive


# --- 1. File Paths and Defaults ---
CONFIG_PATH = '/content/plot_config.json'
DEFAULT_CONFIG = {
    "template": "plotly_white",
    "max_width": 1200,
    "fallback": "#007bff",
    "domain_map": {
        "sent": "#636EFA",
        "click": "#EF553B",
        "control_group": "#00CC96"
    }
}

# --- 2. Config Load/Save ---
def load_config(path=CONFIG_PATH, defaults=DEFAULT_CONFIG):
    if os.path.exists(path):
        try:
            with open(path, 'r') as f:
                return json.load(f)
        except Exception:
            pass
    with open(path, 'w') as f:
        json.dump(defaults, f, indent=2)
    return defaults.copy()

def save_config(cfg):
    with open(CONFIG_PATH, 'w') as f:
        json.dump(cfg, f, indent=2)

def get_plotly_config():
    with open(CONFIG_PATH) as f:
        cfg = json.load(f)
    cfg['domain_map'] = {k.lower(): v for k, v in cfg.get('domain_map', {}).items()}
    return cfg

# --- 3. UI Widgets ---
config = load_config()
template_widget = widgets.Dropdown(
    options=["plotly", "plotly_white", "plotly_dark", "presentation", "simple_white", "ggplot2", "seaborn"],
    value=config.get("template", DEFAULT_CONFIG["template"]),
    description="Template:", style={'description_width': '90px'},
    layout=widgets.Layout(width='290px')
)
width_widget = widgets.BoundedIntText(
    value=config.get("max_width", DEFAULT_CONFIG["max_width"]),
    min=400, max=2400, step=50,
    description="Plot Width:", style={'description_width': '90px'},
    layout=widgets.Layout(width='200px')
)
fallback_widget = widgets.ColorPicker(
    value=config.get("fallback", DEFAULT_CONFIG["fallback"]),
    description="Fallback:", style={'description_width': '90px'},
    layout=widgets.Layout(width='200px')
)

def make_domain_row(domain, color):
    name_widget = widgets.Text(value=domain, layout=widgets.Layout(width='120px'))
    color_widget = widgets.ColorPicker(value=color, layout=widgets.Layout(width='80px'))
    remove_btn = widgets.Button(icon='trash', layout=widgets.Layout(width='32px'))
    row = widgets.HBox([name_widget, color_widget, remove_btn])
    remove_btn.on_click(lambda btn: remove_domain_row(row))
    return row

def remove_domain_row(row):
    with domain_colors_box.hold_trait_notifications():
        children = list(domain_colors_box.children)
        if row in children:
            children.remove(row)
        domain_colors_box.children = tuple(children)

def add_domain_row(_=None):
    new_row = make_domain_row("new_domain", "#222222")
    domain_colors_box.children += (new_row,)

def get_domain_map():
    domain_map = {}
    for row in domain_colors_box.children:
        name = row.children[0].value.strip().lower()
        color = row.children[1].value
        if name:
            domain_map[name] = color
    return domain_map

domain_colors_box = widgets.VBox([
    make_domain_row(dom, col) for dom, col in config.get("domain_map", {}).items()
])
add_domain_btn = widgets.Button(description="Add Domain", icon='plus', button_style='info')
add_domain_btn.on_click(add_domain_row)

out = widgets.Output()
def on_save(_):
    config['template'] = template_widget.value
    config['max_width'] = width_widget.value
    config['fallback'] = fallback_widget.value
    config['domain_map'] = get_domain_map()
    save_config(config)
    with out:
        clear_output()
        print("✅ Config saved!")
    # Update global cfg and Plotly template
    builtins.cfg = get_plotly_config()
    pio.templates.default = builtins.cfg["template"]

save_btn = widgets.Button(description="💾 Save Config", button_style="primary")
save_btn.on_click(on_save)

config_card = widgets.VBox([
    widgets.HTML("<b>Plotly Settings</b>"),
    template_widget,
    width_widget,
    fallback_widget,
    widgets.HTML("<b>Domain Colors</b>"),
    domain_colors_box,
    add_domain_btn,
    save_btn,
    out
], layout=widgets.Layout(
    border="1px solid #e0e0e0", border_radius="10px",
    padding="20px", margin="10px 0", width="370px"
))

display(config_card)

# --- 4. Global config + Plotly template setup ---
builtins.cfg = get_plotly_config()
pio.templates.default = builtins.cfg["template"]

# --- 5. UNIVERSAL REUSABLE PLOT FUNCTION ---
def plotly_group_dropdown(
    df,
    x,
    y_metrics,
    group,
    labels=None,
    kind='bar',
    title_prefix='',
    x_title='',
    legend_title='',
    show=True  # <--- Default is True
):
    import plotly.graph_objects as go
    cfg = builtins.cfg

    # Responsive plot: no static width, use autosize
    fig = go.Figure()
    if isinstance(y_metrics, str):
        metrics = [y_metrics]
    else:
        metrics = y_metrics
    if labels is None:
        labels = {m: m for m in metrics}
    groups = df[group].unique()

    for metric_idx, metric in enumerate(metrics):
        for group_idx, g in enumerate(groups):
            sub = df[df[group] == g]
            trace_args = dict(
                x=sub[x],
                y=sub[metric],
                name=str(g),
                marker_color=cfg["domain_map"].get(str(g).lower(), cfg["fallback"]),
                visible=(metric_idx == 0),
                showlegend=(metric_idx == 0)
            )
            if kind == 'bar':
                fig.add_trace(go.Bar(**trace_args))
            elif kind == 'line':
                fig.add_trace(go.Scatter(mode='lines+markers', **trace_args))

    if len(metrics) > 1:
        n_groups = len(groups)
        n_metrics = len(metrics)
        buttons = []
        for metric_idx, metric in enumerate(metrics):
            visibility = [False] * (n_groups * n_metrics)
            for group_idx in range(n_groups):
                visibility[metric_idx * n_groups + group_idx] = True
            showlegend = [False] * (n_groups * n_metrics)
            for group_idx in range(n_groups):
                if visibility[metric_idx * n_groups + group_idx]:
                    showlegend[metric_idx * n_groups + group_idx] = True
            buttons.append(dict(
                label=labels[metric],
                method="update",
                args=[
                    {"visible": visibility, "showlegend": showlegend},
                    {"title": f"{title_prefix}{labels[metric]}",
                     "yaxis": {"title": labels[metric]}}
                ]
            ))
        fig.update_layout(
            updatemenus=[dict(
                active=0, buttons=buttons, x=1.1, y=0.7, xanchor='center', yanchor='top'
            )],
            title=f"{title_prefix}{labels[metrics[0]]}",
            xaxis_title=x_title,
            yaxis_title=labels[metrics[0]],
            barmode='group' if kind == 'bar' else None,
            legend_title=legend_title,
            template=cfg["template"],
            autosize=True,
            margin=dict(l=20, r=20, t=50, b=20)
        )
    else:
        fig.update_layout(
            title=f"{title_prefix}{labels[metrics[0]]}",
            xaxis_title=x_title,
            yaxis_title=labels[metrics[0]],
            barmode='group' if kind == 'bar' else None,
            legend_title=legend_title,
            template=cfg["template"],
            autosize=True,
            margin=dict(l=20, r=20, t=50, b=20)
        )
    # Make plot responsive
    fig.update_layout(
        autosize=True,
        margin=dict(l=20, r=20, t=50, b=20)
    )
    config_plotly = dict(responsive=True, displayModeBar=True)
    if show:
        fig.show(config=config_plotly)

        # --- PowerPoint Download Button ---
    def on_download_pptx_clicked(b):
        pptx_buffer = dataframe_to_pptx_chart(
            df, x, metrics, group, labels, kind, title_prefix, x_title, legend_title, cfg
        )
        fname = f"{title_prefix.replace(' ', '_').lower()}chart.pptx"
        with open(fname, "wb") as f:
            f.write(pptx_buffer.read())
        # Need to import files for this to work in Colab
        from google.colab import files
        files.download(fname)


    download_pptx_btn = widgets.Button(description="⬇️ Download as PowerPoint")
    download_pptx_btn.on_click(on_download_pptx_clicked)
    display(download_pptx_btn)

        # --- Excel Download Button ---
    def on_download_excel_clicked(b):
        excel_buffer = dataframe_to_excel_with_colors(
            df, x, metrics, group, labels, kind, title_prefix, x_title, legend_title, cfg
        )
        fname = f"{title_prefix.replace(' ', '_').lower()}chart.xlsx"
        with open(fname, "wb") as f:
            f.write(excel_buffer.read())
        # Need to import files for this to work in Colab
        from google.colab import files
        files.download(fname)

    download_excel_btn = widgets.Button(description="⬇️ Download as Excel")
    download_excel_btn.on_click(on_download_excel_clicked)
    display(download_excel_btn)


    return fig

# --- 6. (OPTIONAL) Helper: Query and Categorical Ordering ---
def query_categorical(sql, order_col=None, order=None):
    df = client.query(sql).to_dataframe()
    if order_col and order:
        df[order_col] = pd.Categorical(df[order_col], categories=order, ordered=True)
        df = df.sort_values(by=[order_col])
    return df

def dataframe_to_pptx_chart(
    df, x, y_metrics, group, labels=None, kind='bar', title_prefix='', x_title='', legend_title='', cfg=None
):
    """
    Create a PowerPoint pptx with the same color scheme as Plotly config.
    """
    # Import necessary libraries inside the function
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.dml.color import RGBColor
    import io

    if isinstance(y_metrics, str):
        metrics = [y_metrics]
    else:
        metrics = y_metrics
    if labels is None:
        labels = {m: m for m in metrics}
    if cfg is None:
        cfg = builtins.cfg if hasattr(builtins, "cfg") else get_plotly_config()

    pptx_buffer = io.BytesIO()
    prs = Presentation()

    for metric in metrics:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        chart_data = CategoryChartData()
        categories = list(df[x].unique())
        chart_data.categories = categories
        groups = list(df[group].unique())

        for g in groups:
            sub = df[df[group] == g].set_index(x).reindex(categories)
            yvals = sub[metric].values
            chart_data.add_series(str(g), [0 if pd.isna(v) else v for v in yvals])

        x0, y0, cx, cy = Inches(1), Inches(1.25), Inches(7), Inches(3.5)
        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED if kind == 'bar' else XL_CHART_TYPE.LINE
        chart = slide.shapes.add_chart(chart_type, x0, y0, cx, cy, chart_data).chart

        chart.has_title = True
        chart.chart_title.text_frame.text = f"{title_prefix}{labels[metric]}"
        chart.category_axis.has_title = True
        chart.category_axis.axis_title.text_frame.text = x_title
        chart.value_axis.has_title = True
        chart.value_axis.axis_title.text_frame.text = labels[metric]
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = 2  # bottom

        # --- Apply Plotly Colors! ---
        # Helper function to convert hex to RGBColor for pptx
        def hex_to_rgb_color(hex_color):
            hex_color = hex_color.lstrip("#")
            return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

        for idx, s in enumerate(chart.series):
            group_name = groups[idx]
            col_hex = cfg["domain_map"].get(str(group_name).lower(), cfg["fallback"])
            try:
                s.format.fill.solid()
                s.format.fill.fore_color.rgb = hex_to_rgb_color(col_hex)
            except Exception as e:
                pass  # fallback if color fails

    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer

def hex_to_rgb_tuple(hex_color):
    """Convert hex (#RRGGBB) to (R, G, B) tuple for xlsxwriter."""
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


# --- 5. Universal Excel Export Helper ---
def dataframe_to_excel_with_colors(
    df, x, y_metrics, group=None, labels=None, kind='bar', title_prefix='', x_title='', legend_title='', cfg=None
):
    # Import necessary libraries inside the function
    import io
    import pandas as pd # pandas is already imported, but good practice if function could be standalone
    import xlsxwriter # xlsxwriter is used by pd.ExcelWriter engine

    if isinstance(y_metrics, str):
        metrics = [y_metrics]
    else:
        metrics = y_metrics
    if labels is None:
        labels = {m: m for m in metrics}
    if cfg is None:
        cfg = builtins.cfg if hasattr(builtins, "cfg") else get_plotly_config()
    excel_buffer = io.BytesIO()
    with pd.ExcelWriter(excel_buffer, engine='xlsxwriter') as writer:
        for metric in metrics:
            if group:
                pivot = df.pivot_table(index=x, columns=group, values=metric)
            else:
                pivot = df.set_index(x)[[metric]]
            sheetname = f'{labels[metric]}'
            pivot.to_excel(writer, sheet_name=sheetname)
            worksheet = writer.sheets[sheetname]
            group_names = list(pivot.columns) if group else [metric]
            for col_idx, col_name in enumerate(group_names, start=1):
                col_hex = cfg["domain_map"].get(str(col_name).lower(), cfg["fallback"])
                rgb = hex_to_rgb_tuple(col_hex)
                fmt = writer.book.add_format({'bg_color': '#{:02X}{:02X}{:02X}'.format(*rgb), 'bold': True})
                worksheet.write(0, col_idx, col_name, fmt)
            chart_type = 'column' if kind == 'bar' else 'line'
            chart = writer.book.add_chart({'type': chart_type})
            n_rows = len(pivot)
            for col_idx, col_name in enumerate(group_names):
                col_hex = cfg["domain_map"].get(str(col_name).lower(), cfg["fallback"])
                rgb = hex_to_rgb_tuple(col_hex)
                series = {
                    'name':       [sheetname, 0, col_idx+1],
                    'categories': [sheetname, 1, 0, n_rows, 0],
                    'values':     [sheetname, 1, col_idx+1, n_rows, col_idx+1],
                    'fill':       {'color': '#{:02X}{:02X}{:02X}'.format(*rgb)},
                }
                if chart_type == 'line':
                    series['line'] = {'color': '#{:02X}{:02X}{:02X}'.format(*rgb)}
                chart.add_series(series)
            chart.set_title({'name': f'{labels[metric]} by {x}'})
            chart.set_x_axis({'name': x_title or x})
            chart.set_y_axis({'name': labels[metric]})
            chart.set_legend({'position': 'bottom'})
            worksheet.insert_chart(2, len(group_names) + 3, chart)
    excel_buffer.seek(0)
    return excel_buffer

# --- 6. Register the Excel Export as a Colab Callback for Modebar ---
# Import the necessary module
from google.colab import output, files

def download_excel_callback():
    # Uses latest global variables defined by the user
    # Ensure these variables are accessible or passed correctly if this is a true callback
    # For the purpose of fixing the NameError, assuming they are globally available or handled by Colab's execution context
    global export_df, export_x, export_y_metrics, export_group, export_labels, export_kind, export_title_prefix, export_x_title, export_legend_title, cfg # Assuming these are meant to be global
    excel_buffer = dataframe_to_excel_with_colors(
        export_df, export_x, export_y_metrics, export_group, export_labels,
        kind=export_kind, title_prefix=export_title_prefix, x_title=export_x_title,
        legend_title=export_legend_title, cfg=cfg
    )
    fname = "plotly_chart_export.xlsx"
    with open(fname, "wb") as f:
        f.write(excel_buffer.read())
    files.download(fname)

# Register the callback after importing output
output.register_callback("notebook.download_excel_via_modebar", download_excel_callback)

def add_excel_modebar_button(fig):
    spreadsheet_svg = "M512 64H0v384h512V64zm-32 352H32V96h448v320zM80 144h96v64H80v-64zm128 0h224v64H208v-64zm0 96h224v64H208v-64zm-128 0h96v64H80v-64z"
    custom_btn = {
        "name": "Export Excel",
        "icon": {
            "width": 1000,
            "height": 1000,
            "path": spreadsheet_svg,
            "ascent": 1000, "descent": 0
        },
        "click": "function(gd){google.colab.kernel.invokeFunction('notebook.download_excel_via_modebar', [], {});}"
    }
    config_plotly = dict(
        responsive=True,
        displayModeBar=True,
        modeBarButtonsToAdd=[custom_btn]
    )
    fig.show(config=config_plotly)
    print("🔗 To download Excel: use the green spreadsheet button in the chart's toolbar.")

def display_kpi_dashboard(kpis, title=None, subtitle=None, card_width=170, icon_map=None):
    """
    kpis: list of dicts, each with keys 'label', 'value', and optional 'color', 'icon'
    title: Optional title above the cards
    subtitle: Optional subtitle below title
    card_width: width of each card in px
    icon_map: dict (label -> icon unicode/emoji)
    """
    if icon_map is None:
        icon_map = {}

    cards = []
    for kpi in kpis:
        label = kpi.get('label', '')
        value = kpi.get('value', '')
        color = kpi.get('color', '#1976d2')
        icon = kpi.get('icon') or icon_map.get(label, '')
        icon_html = f"<span style='font-size:22px;margin-right:5px'>{icon}</span>" if icon else ""
        html = f"""
        <div style="
            display: flex; flex-direction: column; align-items: center; justify-content: center;
            background: #fff;
            border-left: 7px solid {color};
            border-radius: 14px;
            box-shadow: 0 4px 14px #0002;
            min-width: {card_width}px; max-width: {card_width+30}px;
            margin: 10px 12px 10px 0;
            padding: 18px 16px 12px 16px;
        ">
          <div style="font-size: 16px; color: #444;">{icon_html}{label}</div>
          <div style="font-size: 32px; font-weight: bold; color: {color}; margin-top:10px;">{value}</div>
        </div>
        """
        cards.append(widgets.HTML(html))
    dashboard = widgets.HBox(cards)
    items = []
    if title:
        # Wrap IPython.display.HTML object in ipywidgets.HTML
        items.append(widgets.HTML(f"<h2 style='color:#1976d2;margin-bottom:2px'>{title}</h2>"))
    if subtitle:
        # Wrap IPython.display.HTML object in ipywidgets.HTML
        items.append(widgets.HTML(f"<div style='color:#888;margin-bottom:8px'>{subtitle}</div>"))
    items.append(dashboard)
    display(widgets.VBox(items))